#!/usr/bin/env python3
"""Practice deck: TA architecture order + original PNGs. White 16:9."""
from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu, Inches

ROOT = Path(__file__).resolve().parents[2]
IMG = ROOT / "06_defense" / "配图"
OUT_DIR = Path(__file__).resolve().parent / "练习预览"
PPTX = ROOT / "06_defense" / "13组-摸鱼-答辩练习.pptx"
FONT = "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"

W, H = 1920, 1080
BG = (255, 255, 255)
INK = (32, 32, 32)
MUTED = (90, 90, 90)
BLUE = (31, 78, 121)
LINE = (210, 216, 224)
BAND = (242, 246, 250)
RED = (192, 80, 77)
WHITE = (255, 255, 255)
TOTAL = 8


def font(n: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT, n)


def new() -> Image.Image:
    return Image.new("RGB", (W, H), BG)


def header(im: Image.Image, title: str, page: int) -> ImageDraw.ImageDraw:
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 8), fill=BLUE)
    d.text((48, 24), title, font=font(32), fill=BLUE)
    d.line((48, 74, W - 48, 74), fill=LINE, width=2)
    d.text((48, H - 40), f"13组 摸鱼  ·  {page}/{TOTAL}", font=font(16), fill=MUTED)
    d.line((48, H - 54, W - 48, H - 54), fill=LINE, width=1)
    return d


def table(d, origin, col_w, row_h, rows, fs=16):
    x0, y0 = origin
    for r, row in enumerate(rows):
        y = y0 + r * row_h
        bg = BLUE if r == 0 else (BAND if r % 2 else WHITE)
        fg = WHITE if r == 0 else INK
        cx = x0
        for i, cell in enumerate(row):
            d.rectangle((cx, y, cx + col_w[i], y + row_h), fill=bg, outline=LINE)
            use = font(fs)
            s = str(cell)
            pad = 8
            max_w = col_w[i] - 2 * pad
            lines = wrap(d, s, use, max_w)
            while (len(lines) * (use.size + 4) > row_h - 4 or any(d.textlength(ln, font=use) > max_w for ln in lines)) and use.size > 11:
                use = font(use.size - 1)
                lines = wrap(d, s, use, max_w)
            block_h = len(lines) * (use.size + 3)
            ty = y + max(2, (row_h - block_h) // 2)
            for ln in lines:
                d.text((cx + pad, ty), ln, font=use, fill=fg)
                ty += use.size + 3
            cx += col_w[i]


def paste(im, name, box):
    path = IMG / name
    x, y, w, h = box
    src = Image.open(path).convert("RGB")
    scale = min(w / src.width, h / src.height)
    nw, nh = max(1, int(src.width * scale)), max(1, int(src.height * scale))
    src = src.resize((nw, nh), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (w, h), WHITE)
    canvas.paste(src, ((w - nw) // 2, (h - nh) // 2))
    ImageDraw.Draw(im).rectangle((x - 1, y - 1, x + w, y + h), outline=LINE)
    im.paste(canvas, (x, y))


def wrap(d, text, f, max_w):
    lines, cur = [], ""
    for ch in text:
        t = cur + ch
        if d.textlength(t, font=f) <= max_w:
            cur = t
        else:
            if cur:
                lines.append(cur)
            cur = ch
    lines.append(cur)
    return lines


SLIDES: list[Image.Image] = []
NOTES: list[str] = []


def add(im, note):
    SLIDES.append(im)
    NOTES.append(note)


def s1():
    im = new()
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 16, H), fill=BLUE)
    d.rectangle((0, 0, W, 8), fill=BLUE)
    d.text((56, 200), "软件工程基础实践  2026 夏", font=font(26), fill=MUTED)
    d.text((56, 260), "摸鱼", font=font(88), fill=BLUE)
    d.text((56, 380), "校园购物 + 二手交易平台", font=font(36), fill=INK)
    d.line((56, 450, 520, 450), fill=BLUE, width=3)
    d.text((56, 480), "杨任宇老师班  ·  13组", font=font(28), fill=INK)
    d.text((56, 560), "https://github.com/tchen-0213/softw", font=font(22), fill=BLUE)
    d.text((56, 640), "鲁在精  浦灵一  王悠然  赵紫嫣  陈子正  剧博洋", font=font(22), fill=INK)
    paste(im, "P00-仓库首页.png", (1020, 160, 820, 760))
    add(im, "报组号、项目名、仓库。")


def s2():
    im = new()
    d = header(im, "项目目标及全部业务场景完成情况", 2)
    d.text((48, 90), "校园购物 + 二手。原系统 React + Express + MySQL shopping_platform，标签 monolith-start（10fa639）", font=font(18), fill=INK)
    d.text((48, 122), "小学期完成 Docker 三容器、Actions、K8s，以及 user / product / order 与网关，标签 microservices-v1（63585e0）", font=font(18), fill=INK)
    d.text((48, 154), "单体 8080 / 3001；微服务 8082 / 8081。UC01–UC12 均已完成。", font=font(18), fill=INK)
    rows = [["编号", "场景", "状态"]]
    data = [
        ("UC01", "注册并登录", "已完成"), ("UC02", "浏览并搜索", "已完成"), ("UC03", "详情加购", "已完成"),
        ("UC04", "下单支付发货收货", "已完成"), ("UC05", "发布二手", "已完成"), ("UC06", "店铺与商品", "已完成"),
        ("UC07", "评价", "已完成"), ("UC08", "聊天议价", "已完成"), ("UC09", "地址", "已完成"),
        ("UC10", "订单与物流", "已完成"), ("UC11", "取消并恢复库存", "已完成"), ("UC12", "公开店铺与信用", "已完成"),
    ]
    for a, b, c in data:
        rows.append([a, b, c])
    table(d, (48, 195), [110, 420, 120], 52, rows, 17)
    paste(im, "P01-总用例图.png", (740, 195, 1132, 780))
    add(im, "12 个都完成。现场重点三条。")


def s3():
    im = new()
    d = header(im, "代表性用例：需求 — 设计 — 代码 — 测试", 3)
    table(
        d, (48, 90), [100, 170, 180, 170, 620, 540], 52,
        [
            ["用例", "系统级", "组件级", "对象级", "代码", "测试"],
            ["UC01", "SYS-SEQ01", "COMP-SEQ01", "OBJ-SEQ01", "user.js / userController / AuthPage", "UNIT/INT/E2E-TC01"],
            ["UC02", "SYS-SEQ02", "COMP-SEQ02", "OBJ-SEQ02", "product.js / productController / SearchPage", "UNIT/INT/E2E-TC02"],
            ["UC04", "SYS-SEQ04", "COMP-SEQ04", "OBJ-SEQ04", "order.js / orderController / CheckoutPage", "UNIT/INT/E2E-TC04"],
        ],
        16,
    )
    d.text((48, 320), "系统级顺序图：游客 / 买家 / 卖家 与平台的交互", font=font(18), fill=BLUE)
    paste(im, "P01-SYS-SEQ01-注册登录.png", (48, 355, 600, 620))
    paste(im, "P01-SYS-SEQ02-搜索.png", (668, 355, 600, 620))
    paste(im, "P01-SYS-SEQ04-下单履约.png", (1288, 355, 584, 620))
    add(im, "点 UC04：需求图、代码 orderController、测试 E2E-TC04 能对上。")


def s4():
    im = new()
    d = header(im, "模型到代码（以 UC04 为例）", 4)
    table(
        d, (48, 90), [280, 520, 420, 560], 48,
        [
            ["层次", "模型", "代码", "图"],
            ["系统级 SYS-SEQ04", "买家、卖家；下单、支付、发货、收货", "下单到收货闭环", "SYS-SEQ04"],
            ["组件级 COMP-SEQ04", "结账页、鉴权、订单控制器、事务", "orderController、CheckoutPage", "COMP-SEQ04"],
            ["对象级 OBJ-SEQ04", "createOrder()、pay / ship / confirm", "与源码函数名一致", "OBJ-SEQ04"],
        ],
        16,
    )
    d.text((48, 300), "状态：待付款 → 待发货 → 待收货 → 已完成。order-service 按 reservationId 调 product-service 预留/释放/完成。", font=font(18), fill=INK)
    paste(im, "P01-COMP-SEQ04-下单组件.png", (48, 345, 620, 630))
    paste(im, "P01-OBJ-SEQ04-下单对象.png", (688, 345, 600, 630))
    paste(im, "P04-订单状态机.png", (1308, 345, 564, 630))
    add(im, "三层图指到 createOrder 即可，不要念每条箭头。")


def s5():
    im = new()
    d = header(im, "业务微服务：职责、划分依据、接口、表归属", 5)
    d.text((48, 88), "调用链：前端 → API 网关 → user-service / product-service / order-service。按身份、可交易资源、交易过程划分。", font=font(18), fill=INK)
    table(
        d, (48, 125), [240, 520, 360, 660], 44,
        [
            ["服务", "职责", "划分依据", "库与表"],
            ["user-service", "注册登录、资料、密码、地址、角色、信用", "身份与账号生命周期", "softw_users：Users, Addresses"],
            ["product-service", "商品、二手、店铺、评价、聊天、议价、图片、库存", "可交易资源", "softw_catalog + uploads"],
            ["order-service", "下单、支付、取消、发货、收货、状态机", "交易过程与订单快照", "softw_orders：Orders, OrderSellers"],
            ["API 网关", "路由、JWT 透传、超时、CORS", "统一入口", "无库"],
        ],
        15,
    )
    d.text((48, 360), "网关路由：/api/users、/api/addresses → user；商品、二手、店铺、评价、聊天、上传 → product；/api/orders → order。一表由一个服务写；订单保存快照；内部调用带 X-Internal-Token。", font=font(16), fill=INK)
    paste(im, "P03-微服务划分.png", (48, 400, 1800, 570))
    add(im, "三个服务 + 网关不算第四个。评价聊天在商品服务。")


def s6():
    im = new()
    d = header(im, "调用失败时的处理办法", 6)
    table(
        d, (48, 100), [480, 1340], 78,
        [
            ["场景", "处理"],
            ["网关超时或上游断开", "返回 HTTP 503"],
            ["库存预留失败（无货或超时）", "不创建订单；超时返回 503"],
            ["释放 / 完成预留", "按 reservationId 幂等，重复调用不重复加库存"],
            ["停 product-service", "商品接口 503「依赖服务暂不可用」；订单依赖检查 206 degraded「商品信息暂不可用，订单查询保持可用」；网关、用户、订单仍 1/1"],
            ["内部库存、信用、购买证明", "请求须带 X-Internal-Token"],
        ],
        18,
    )
    d.text((48, 560), "实验：npm run experiment:fault。恢复后商品接口 200，副本重建，无残留。", font=font(22), fill=INK)
    add(im, "抓 503 和 206 两张牌。")


def s7():
    im = new()
    d = header(im, "关键数据：测试 · CI/CD · HPA · 故障 · 性能", 7)
    table(
        d, (48, 88), [150, 880], 42,
        [
            ["层级", "结果（2026-09-03）"],
            ["单元", "后端 220 通过 / 0 失败 / 1 跳过（含 REG-BE 100/100）；前端原有 100/100，新增 REG-FE 已注册"],
            ["覆盖率", "语句 94.42% / 分支 81.83% / 函数 92.34% / 行 94.42%"],
            ["集成", "单体 API 32/32；隔离 22/22（网关 17/17）；网关 API/E2E 15/15；公开 API 49 项"],
            ["端到端", "Playwright 单体与微服务均为 42/42（12 + 27 + 3）"],
            ["CI/CD", "工作流 softw-ci-cd：测试通过后构建 7 个镜像并部署 Kind；测试失败不构建。Actions #77（33579985248）"],
            ["HPA", "1 → 3 → 5 → 1；997 请求、9.49 req/s、错误 0/997"],
            ["故障", "停商品服务后：503 / 206；网关、用户、订单仍存活"],
        ],
        15,
    )
    d.text((48, 440), "单体与微服务对比（2026-08-28，同机、同商品数据、同一 k6，18 组，错误率 0%）", font=font(16), fill=INK)
    table(
        d, (48, 478), [220, 200, 200, 180, 180], 40,
        [
            ["接口", "吞吐", "延迟", "CPU", "内存"],
            ["商品列表", "+61.9%", "−38.6%", "+55.2%", "+5.1%"],
            ["商品搜索", "+53.9%", "−35.5%", "+50.7%", "+17.3%"],
            ["商品详情", "+30.9%", "−26.2%", "+6.9%", "+14.0%"],
        ],
        16,
    )
    paste(im, "P02-GitHub-Actions-77.png", (1080, 88, 792, 430))
    paste(im, "P02-流水线失败阻断-20260825.png", (1080, 540, 792, 430))
    d.text((1080, 978), "上：Actions #77    下：2026-08-25 测试失败，未进入部署", font=font(14), fill=MUTED)
    add(im, "报 220、42/42、#77、1到5再回1、+61.9%。再说不是天然更快。")


def s8():
    im = new()
    d = header(im, "成员分工", 8)
    table(
        d, (48, 90), [280, 720], 52,
        [
            ["成员", "职责"],
            ["鲁在精", "项目管理；商品搜索与浏览"],
            ["浦灵一", "在线下单与支付；后端与安全"],
            ["王悠然", "物流跟踪、订单状态"],
            ["赵紫嫣", "店铺管理；测试"],
            ["陈子正", "信用评价；流水线与部署"],
            ["剧博洋", "二手交易；性能与数据库"],
        ],
        18,
    )
    d.text((48, 480), "大模型用于仓库检索、测试建议、文档整理和部署检查。", font=font(18), fill=INK)
    d.text((48, 524), "代码阅读、测试执行、版本核对和凭据扫描由组员完成。", font=font(18), fill=INK)
    paste(im, "P09-GitHub-Pages演示.png", (1080, 90, 792, 880))
    add(im, "分工一句。切录屏。")


def pack(pptx_path: Path):
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, im in enumerate(SLIDES):
        png = OUT_DIR / f"S{i + 1:02d}.png"
        im.save(png, "PNG", optimize=True)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = NOTES[i]
    prs.save(pptx_path)


def main():
    for fn in (s1, s2, s3, s4, s5, s6, s7, s8):
        fn()
    pack(PPTX)
    print(PPTX)


if __name__ == "__main__":
    main()
