#!/usr/bin/env python3
"""White, dense course-defense slides. Numbers copied from submitted reports."""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu, Inches

ROOT = Path(__file__).resolve().parents[2]
PNG = ROOT / "02_docs" / "images" / "png"
PIPE = ROOT / "05_management" / "流水线截图"
OUT_DIR = Path(__file__).resolve().parent / "slides"
PPTX = ROOT / "06_defense" / "13组-摸鱼-最终答辩.pptx"
FONT = "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"

W, H = 1920, 1080
BG = (255, 255, 255)
INK = (32, 32, 32)
MUTED = (96, 96, 96)
BLUE = (31, 78, 121)
BLUE2 = (47, 117, 181)
LINE = (210, 216, 224)
BAND = (242, 246, 250)
RED = (192, 80, 77)
GREEN = (84, 130, 53)
WHITE = (255, 255, 255)


def font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT, size)


def new_slide() -> Image.Image:
    return Image.new("RGB", (W, H), BG)


def draw_header(im: Image.Image, title: str, page: int, total: int) -> ImageDraw.ImageDraw:
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 8), fill=BLUE)
    d.text((56, 28), title, font=font(36), fill=BLUE)
    d.line((56, 82, W - 56, 82), fill=LINE, width=2)
    foot = f"杨任宇老师班 · 13组  摸鱼    {page} / {total}"
    d.text((56, H - 42), foot, font=font(16), fill=MUTED)
    d.line((56, H - 56, W - 56, H - 56), fill=LINE, width=1)
    return d


def text_w(d: ImageDraw.ImageDraw, text: str, f: ImageFont.FreeTypeFont) -> int:
    return int(d.textlength(text, font=f))


def wrap(d: ImageDraw.ImageDraw, text: str, f: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
    lines: list[str] = []
    for para in text.split("\n"):
        if para == "":
            lines.append("")
            continue
        cur = ""
        for ch in para:
            trial = cur + ch
            if d.textlength(trial, font=f) <= max_w:
                cur = trial
            else:
                if cur:
                    lines.append(cur)
                cur = ch
        lines.append(cur)
    return lines


def draw_wrapped(d: ImageDraw.ImageDraw, text: str, xy: tuple[int, int], f, fill, max_w: int, lh: int) -> int:
    x, y = xy
    for line in wrap(d, text, f, max_w):
        d.text((x, y), line, font=f, fill=fill)
        y += lh
    return y


def table(d: ImageDraw.ImageDraw, origin, col_w, row_h, rows, header=True, font_size=18):
    x0, y0 = origin
    f = font(font_size)
    fh = font(font_size)
    for r, row in enumerate(rows):
        y = y0 + r * row_h
        bg = BLUE if r == 0 and header else (BAND if r % 2 else WHITE)
        fg = WHITE if r == 0 and header else INK
        d.rectangle((x0, y, x0 + sum(col_w), y + row_h), fill=bg, outline=LINE)
        cx = x0
        for i, cell in enumerate(row):
            d.rectangle((cx, y, cx + col_w[i], y + row_h), outline=LINE)
            pad = 10
            # shrink if needed
            use = f if r else fh
            s = str(cell)
            while d.textlength(s, font=use) > col_w[i] - 2 * pad and use.size > 12:
                use = font(use.size - 1)
            d.text((cx + pad, y + (row_h - use.size) // 2 - 1), s, font=use, fill=fg)
            cx += col_w[i]


def paste_contain(im: Image.Image, path: Path, box, bg=WHITE):
    x, y, w, h = box
    src = Image.open(path).convert("RGB")
    scale = min(w / src.width, h / src.height)
    nw, nh = max(1, int(src.width * scale)), max(1, int(src.height * scale))
    src = src.resize((nw, nh), Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", (w, h), bg)
    canvas.paste(src, ((w - nw) // 2, (h - nh) // 2))
    d = ImageDraw.Draw(im)
    d.rectangle((x - 1, y - 1, x + w, y + h), outline=LINE, width=1)
    im.paste(canvas, (x, y))


def panel(d, box, fill=BAND):
    x, y, w, h = box
    d.rectangle((x, y, x + w, y + h), fill=fill, outline=LINE)


NOTES: list[str] = []
SLIDES: list[Image.Image] = []
TOTAL = 18


def add(im: Image.Image, note: str):
    SLIDES.append(im)
    NOTES.append(note)


def s01():
    im = new_slide()
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 18, H), fill=BLUE)
    d.rectangle((0, 0, W, 8), fill=BLUE)
    d.text((80, 220), "软件工程基础实践  2026 夏", font=font(28), fill=MUTED)
    d.text((80, 280), "摸鱼", font=font(92), fill=BLUE)
    d.text((80, 400), "校园购物 + 二手交易平台", font=font(40), fill=INK)
    d.line((80, 470, 520, 470), fill=BLUE2, width=3)
    d.text((80, 500), "杨任宇老师班  ·  13组  ·  最终答辩", font=font(28), fill=INK)
    d.text((80, 560), "2026-09-04    10 分钟（架构 3′ / 演示 4′ / 问答 3′）", font=font(22), fill=MUTED)
    d.text((80, 640), "https://github.com/tchen-0213/softw", font=font(22), fill=BLUE2)
    d.text((80, 720), "鲁在精  浦灵一  王悠然  赵紫嫣  陈子正  剧博洋", font=font(22), fill=INK)
    paste_contain(im, PNG / "01-REQ-USECASE.png", (1020, 160, 820, 760))
    d = ImageDraw.Draw(im)
    d.text((1020, 930), "图  总用例 UC01–UC12", font=font(16), fill=MUTED)
    add(im, "报课程、13组、摸鱼、仓库。不要在封面上念技术细节。")


def s02():
    im = new_slide()
    d = draw_header(im, "答辩 10 分钟怎么用", 2, TOTAL)
    table(
        d,
        (56, 110),
        [180, 160, 1450],
        70,
        [
            ["段落", "时长", "老师要求讲清 / 看到的东西"],
            ["项目和架构", "3 分钟", "目标、全部业务场景、代表用例追溯、模型到代码、三服务职责/划分/接口/表/失败处理、测试·CI·HPA·故障·性能关键数据"],
            ["现场演示", "4 分钟", "提交并触发流水线；K8s/Pod；代表业务流程+自动化测试；HPA 从触发到回落；停一个依赖服务"],
            ["问答", "3 分钟", "组内分工、大模型用什么/做什么/怎么人工检查"],
        ],
        font_size=20,
    )
    panel(d, (56, 420, 1808, 560))
    y = 445
    d.text((80, y), "今天对表要记住的口径", font=font(24), fill=BLUE)
    bullets = [
        "不演示全部用例，核心场景即可：UC01 注册登录、UC02 搜索、UC04 下单支付发货收货；范围仍是 UC01–UC12 全部完成。",
        "拆成微服务后还要测试，交的是第二阶段结果，不是第一阶段。",
        "扩缩容要看到压力升高 Pod 增加、下降后回落；可加速剪辑；不要放进 CI/CD。",
        "故障：停一个微服务后，其它业务还能用，系统给出事先设计的 503 / 206 提示或备用结果。",
        "录屏最好内嵌 PPT，配音+字幕；现场以可重复真实过程为准。",
    ]
    y = 500
    for b in bullets:
        y = draw_wrapped(d, "·  " + b, (80, y), font(22), INK, 1740, 34)
        y += 12
    add(im, "先把 3+4+3 报清楚，再进入项目。")


def s03():
    im = new_slide()
    d = draw_header(im, "项目目标与小学期改造", 3, TOTAL)
    panel(d, (56, 110, 900, 860))
    d.text((80, 130), "要做成什么", font=font(24), fill=BLUE)
    lines = [
        "校园购物 + 二手：注册登录、检索加购、下单支付发货收货、二手发布、店铺、评价、聊天议价、地址。",
        "原系统：React/Vite 前端、一个 Express、Sequelize、MySQL 库 shopping_platform。",
        "基线标签 monolith-start（10fa639），改造中不移动该标签。",
        "小学期四件套：① Docker 三容器 ② GitHub Actions ③ Kubernetes YAML（探针、HPA）④ 三个业务微服务 + API 网关。",
        "改造后标签 microservices-v1（63585e0）。",
        "单体端口 前端 8080 / 后端 3001；微服务前端 8082、网关 8081。",
        "公网：GitHub Pages 前端 + Codespaces 后端。两套库，登录态和订单不同步。",
    ]
    y = 180
    for t in lines:
        y = draw_wrapped(d, "· " + t, (80, y), font(21), INK, 850, 32)
        y += 10
    paste_contain(im, PNG / "12-ARCH-COMP.png", (990, 110, 870, 860))
    d = ImageDraw.Draw(im)
    d.text((990, 980), "图  原系统组件关系（ARCH-COMP）", font=font(16), fill=MUTED)
    add(im, "原系统一个后端一个库；小学期加上容器、流水线、K8s 和三服务。")


def s04():
    im = new_slide()
    d = draw_header(im, "全部业务场景 UC01–UC12（均已完成）", 4, TOTAL)
    rows = [
        ["编号", "业务场景", "可验证结果", "汇报"],
        ["UC01", "用户注册并登录", "返回 token，可进个人中心", "重点"],
        ["UC02", "浏览并搜索商品", "按条件展示列表", "重点"],
        ["UC03", "详情并加入购物车", "库存/卖家/评价，购物车变化", "覆盖"],
        ["UC04", "创建订单并支付发货收货", "订单状态变更，库存扣减", "重点"],
        ["UC05", "发布二手商品", "进入二手列表可被搜索", "覆盖"],
        ["UC06", "卖家管理店铺和商品", "资料和商品保存成功", "覆盖"],
        ["UC07", "评价已完成订单", "评价可在商品侧查询", "覆盖"],
        ["UC08", "聊天议价并按成交价购买", "一次性核销，他人不可用", "覆盖"],
        ["UC09", "维护收货地址", "结算页可选", "覆盖"],
        ["UC10", "查询本人订单与物流", "状态与轨迹完整", "覆盖"],
        ["UC11", "取消订单并恢复库存", "已取消且库存恢复", "覆盖"],
        ["UC12", "公开店铺与卖家信用", "店铺、信用、在售可查", "覆盖"],
    ]
    table(d, (56, 105), [120, 420, 520, 120], 52, rows, font_size=18)
    paste_contain(im, PNG / "01-REQ-USECASE.png", (1280, 105, 584, 780))
    d = ImageDraw.Draw(im)
    d.text((56, 800), "中期口头范围常为 UC01–UC09；UC10–UC12 由已实现能力在终期纳入追溯。重点只讲三条，不等于只做了三条。", font=font(20), fill=INK)
    add(im, "12 个都做完。今天现场只走 UC01、UC02、UC04。")


def s05():
    im = new_slide()
    d = draw_header(im, "代表用例追溯：需求 → 三级模型 → 代码 → 测试", 5, TOTAL)
    table(
        d,
        (56, 105),
        [110, 170, 180, 170, 520, 480],
        58,
        [
            ["用例", "系统级", "组件级", "对象级", "代码", "测试"],
            ["UC01", "SYS-SEQ01", "COMP-SEQ01", "OBJ-SEQ01", "user.js / userController / AuthPage", "UNIT/INT/E2E-TC01"],
            ["UC02", "SYS-SEQ02", "COMP-SEQ02", "OBJ-SEQ02", "product.js / productController / SearchPage", "UNIT/INT/E2E-TC02"],
            ["UC04", "SYS-SEQ04", "COMP-SEQ04", "OBJ-SEQ04", "order.js / orderController / CheckoutPage", "UNIT/INT/E2E-TC04"],
        ],
        font_size=18,
    )
    d.text((56, 360), "系统级顺序图（参与者 + 平台黑盒，不出现 Controller）", font=font(20), fill=BLUE)
    paste_contain(im, PNG / "03-SYS-SEQ01.png", (56, 400, 580, 520))
    paste_contain(im, PNG / "04-SYS-SEQ02.png", (670, 400, 580, 520))
    paste_contain(im, PNG / "06-SYS-SEQ04.png", (1284, 400, 580, 520))
    d = ImageDraw.Draw(im)
    d.text((56, 930), "左 UC01 注册登录 · 中 UC02 检索 · 右 UC04 订单闭环。完整 12 条见业务场景用例清单与追溯表。", font=font(18), fill=MUTED)
    add(im, "点 UC04：需求图画人怎么走，代码在 orderController，测试编号能对上。")


def s06():
    im = new_slide()
    d = draw_header(im, "模型到代码（UC04 下单支付发货收货）", 6, TOTAL)
    panel(d, (56, 105, 600, 860))
    d.text((76, 125), "三层分别对应什么", font=font(22), fill=BLUE)
    y = 175
    blocks = [
        "需求 SYS-SEQ04：买家下单支付，卖家发货，买家收货。",
        "概要 COMP-SEQ04：Checkout/OrderPage → 鉴权 → 订单控制器 → 商品/订单模型与事务。",
        "详细 OBJ-SEQ04：createOrder() 扣库存；pay / ship / confirm。",
        "状态：待付款 → 待发货 → 待收货 → 已完成；可取消并恢复库存。",
        "微服务：order-service 调 product-service 预留 / 释放 / 完成，按 reservationId 幂等。",
        "测试：UNIT-TC04、INT-TC04-MAIN/ALT/ERR、E2E-TC04。",
    ]
    for t in blocks:
        y = draw_wrapped(d, "· " + t, (76, y), font(20), INK, 560, 30)
        y += 16
    paste_contain(im, PNG / "16-COMP-SEQ04.png", (680, 105, 600, 400))
    paste_contain(im, PNG / "26-OBJ-SEQ04.png", (1300, 105, 564, 400))
    paste_contain(im, PNG / "32-ORDER-STATE.png", (680, 530, 1184, 435))
    d = ImageDraw.Draw(im)
    d.text((680, 980), "上左 组件级 · 上右 对象级 · 下 订单状态机", font=font(16), fill=MUTED)
    add(im, "不要在这页念每一条箭头，指状态机和 createOrder 即可。")


def s07():
    im = new_slide()
    d = draw_header(im, "三个业务微服务：职责、划分、数据", 7, TOTAL)
    table(
        d,
        (56, 100),
        [280, 720, 640],
        52,
        [
            ["服务", "职责", "划分依据"],
            ["user-service", "注册登录、资料、密码、地址、角色、信用  → 库 softw_users", "身份 / 账号生命周期"],
            ["product-service", "商品、二手、店铺、评价、聊天、议价、图片、库存  → softw_catalog + uploads", "可交易资源（标的）"],
            ["order-service", "下单、支付、取消、发货、收货、状态机  → 库 softw_orders", "交易过程与历史快照"],
            ["API Gateway", "路由、JWT 透传、超时、CORS、失败 503；无库", "统一入口，不计入 3 个业务服务"],
        ],
        font_size=18,
    )
    d.text((56, 385), "前端 → API Gateway → user / product / order。不是 12 个用例拆 12 个服务。评价、聊天、上传在 product-service。", font=font(20), fill=INK)
    paste_contain(im, PNG / "33-MS-SPLIT.png", (56, 430, 1808, 530))
    add(im, "身份、标的、交易过程。网关不算第四个业务服务。")


def s08():
    im = new_slide()
    d = draw_header(im, "接口归属、表归属、调用失败怎么处理", 8, TOTAL)
    d.text((56, 100), "网关路径", font=font(22), fill=BLUE)
    table(
        d,
        (56, 135),
        [820, 980],
        40,
        [
            ["前缀", "目标"],
            ["/api/users  /api/addresses", "user-service"],
            ["/api/products  /api/secondhand  /api/shops", "product-service"],
            ["/api/evaluations  /api/chats  /api/uploads  /uploads", "product-service"],
            ["/api/orders", "order-service"],
        ],
        font_size=18,
    )
    d.text((56, 360), "每张表只由一个服务写入；订单保存商品项和地址快照，禁止跨库联表", font=font(22), fill=BLUE)
    table(
        d,
        (56, 400),
        [280, 900, 620],
        40,
        [
            ["库 / 卷", "表", "唯一管理"],
            ["softw_users", "Users, Addresses", "user-service"],
            ["softw_catalog", "Products, Shops, Evaluations, ChatConversations, ChatMessages, InventoryReservations", "product-service"],
            ["softw_orders", "Orders", "order-service"],
            ["uploads/", "图片文件", "product-service"],
        ],
        font_size=17,
    )
    d.text((56, 630), "失败处理（老师会问）", font=font(22), fill=BLUE)
    table(
        d,
        (56, 670),
        [620, 1180],
        48,
        [
            ["调用 / 场景", "处理"],
            ["网关超时或上游断开", "HTTP 503，不伪造成功业务结果"],
            ["库存预留失败（无货/超时）", "不创建订单；超时 503"],
            ["释放 / 完成预留", "按 reservationId 幂等，重复调用不重复加库存"],
            ["商品服务被停掉", "商品 503；订单依赖检查 206 + 备用提示；网关/用户/订单仍 1/1"],
            ["内部库存、信用、购买证明", "必须 X-Internal-Token；外部客户端不能直打"],
        ],
        font_size=17,
    )
    add(im, "失败处理抓 503 和 206 两张牌。内部接口有 token。")


def s09():
    im = new_slide()
    d = draw_header(im, "测试关键数据（2026-09-03，第二阶段）", 9, TOTAL)
    table(
        d,
        (56, 105),
        [480, 1320],
        56,
        [
            ["层级", "结果（原样抄自测试报告，不口算）"],
            ["后端单元 / 安全 / 控制器", "220 通过，0 失败，1 个 API 父入口按设计跳过；其中新增回归 REG-BE-001～100 为 100/100"],
            ["前端 Vitest", "原有基线 100/100；新增 REG-FE-001～100 已注册，Vitest 实测待依赖恢复，不能写成已经全过"],
            ["前端覆盖率", "语句 94.42%  分支 81.83%  函数 92.34%  行 94.42%（门禁 ≥80%）"],
            ["单体真实 MySQL API", "32/32，覆盖 UC01–UC12"],
            ["Playwright E2E", "单体与微服务入口均 42/42（12 主流程 + 27 备选/异常 + 3 组合）"],
            ["微服务静态 / 公共层", "18/18"],
            ["隔离集成", "22/22（其中网关 17/17）"],
            ["网关 API / E2E", "15/15；公开业务 API 49 项有映射"],
            ["交付与实验证据门禁", "18/18"],
        ],
        font_size=18,
    )
    d.text((56, 740), "不要把微服务三类测试和单体结果加总成一个「自动化总数」。老师问「拆完还测吗」：要，交二阶段。", font=font(22), fill=RED)
    d.text((56, 800), "重点三条的测试编号：E2E-TC01 / E2E-TC02 / E2E-TC04，均可 --grep 单独跑，不依赖执行顺序。", font=font(20), fill=INK)
    add(im, "报几个大数：220、42/42、32/32。前端新增 100 要诚实说待复跑。")


def s10():
    im = new_slide()
    d = draw_header(im, "CI/CD：先测后发，失败不能进镜像", 10, TOTAL)
    panel(d, (56, 110, 820, 860))
    d.text((80, 130), "工作流 softw-ci-cd", font=font(24), fill=BLUE)
    y = 185
    for t in [
        "文件：.github/workflows/ci-cd.yml",
        "顺序：后端测试 → API → 前端单测与构建 → Playwright → 微服务测试 → 构建 7 个版本化镜像 → Kind 部署与健康检查。",
        "门禁：测试失败不构建镜像、不部署。",
        "终期已验证绿勾：GitHub Actions #77，run 33579985248。",
        "右图是 2026-08-25 早期成功形态（当时 #9），用来说明「提交就会跑」。终期请打开 #77，不要把右图说成 #77。",
        "失败样例见同目录 2026-08-25-06-失败运行6.png，证明红叉会挡住后续 Job。",
    ]:
        y = draw_wrapped(d, "· " + t, (80, y), font(21), INK, 760, 32)
        y += 12
    paste_contain(im, PIPE / "2026-08-25-01-流水线整体成功.png", (900, 110, 964, 430))
    paste_contain(im, PIPE / "2026-08-25-06-失败运行6.png", (900, 560, 964, 410))
    d = ImageDraw.Draw(im)
    d.text((900, 980), "上：早期成功 DAG（标注 8/25）  下：失败阻断", font=font(16), fill=MUTED)
    add(im, "演示优先打开网页上的 #77。强调测试失败不能发版。")


def s11():
    im = new_slide()
    d = draw_header(im, "性能对比（2026-08-28，18 组，错误率 0%）", 11, TOTAL)
    d.text((56, 100), "同机 Apple M5、同 Docker、同 200 条商品、同一 k6、每接口 3 次。微服务相对单体：", font=font(20), fill=INK)
    table(
        d,
        (56, 145),
        [280, 320, 320, 320, 320],
        48,
        [
            ["接口", "吞吐", "平均延迟", "CPU", "内存"],
            ["商品列表", "+61.9%", "−38.6%", "+55.2%", "+5.1%"],
            ["商品搜索", "+53.9%", "−35.5%", "+50.7%", "+17.3%"],
            ["商品详情", "+30.9%", "−26.2%", "+6.9%", "+14.0%"],
        ],
        font_size=20,
    )
    d.text((56, 360), "口播用列表三个数：吞吐 +61.9%，延迟 −38.6%，CPU +55.2%。三轮均值如下（单位 req/s、ms、%、MiB）：", font=font(20), fill=INK)
    table(
        d,
        (56, 405),
        [200, 200, 250, 230, 230, 280, 250],
        42,
        [
            ["版本", "接口", "吞吐", "平均", "P95", "平均 CPU", "平均内存"],
            ["单体", "列表", "1015.108", "4.841", "6.285", "134.314%", "540.902"],
            ["微服务", "列表", "1642.989", "2.970", "4.075", "208.504%", "568.751"],
            ["单体", "搜索", "1843.584", "2.674", "3.542", "164.212%", "517.788"],
            ["微服务", "搜索", "2836.447", "1.726", "2.542", "247.414%", "607.429"],
            ["单体", "详情", "984.646", "5.197", "6.793", "129.656%", "511.013"],
            ["微服务", "详情", "1289.072", "3.833", "5.282", "138.607%", "582.778"],
        ],
        font_size=17,
    )
    panel(d, (56, 730, 1808, 240), fill=(252, 242, 242))
    d.text((80, 755), "不能写成「微服务架构天然更快」。", font=font(26), fill=RED)
    d.text((80, 805), "列表更快，和当前实现有关：单体列表会联 Users 拼卖家，微服务商品表带卖家快照。资源更贵也是事实：网关+商品服务两个 Node，再加独立 MySQL。", font=font(20), fill=INK)
    d.text((80, 855), "这是单机 Docker、5 VU、20 秒短实验，不代表公网或多节点 K8s。", font=font(20), fill=MUTED)
    add(im, "三个百分比脱口而出，紧接着说「不是架构天然更快」。")


def s12():
    im = new_slide()
    d = draw_header(im, "HPA：压力升高 → Pod 增加 → 下降回落（不放进 CI/CD）", 12, TOTAL)
    d.text((56, 105), "命令 npm run experiment:hpa    Kind    CPU 目标 60%    min=1 max=5    实测 1 → 3 → 5 → 1（5 个均 Ready）", font=font(20), fill=INK)
    # stepper
    xs = [120, 560, 1000, 1440]
    labels = [("1", "11:17:28 基线", "1/1/1/1"), ("3", "11:18:05 升上来", "3/3/3/3"), ("5", "11:18:23 封顶 Ready", "5/5/5/5"), ("1", "11:19:47 回落", "1/1/1/1")]
    for i, (n, t, p) in enumerate(labels):
        x = xs[i]
        d.ellipse((x, 165, x + 88, 253), outline=BLUE, width=4)
        d.text((x + (28 if n == "1" else 24), 180), n, font=font(40), fill=BLUE)
        d.text((x, 270), t, font=font(20), fill=INK)
        d.text((x, 305), p, font=font(18), fill=MUTED)
        if i < 3:
            d.line((x + 100, 209, xs[i + 1] - 12, 209), fill=BLUE2, width=3)
    table(
        d,
        (56, 360),
        [280, 400, 400, 720],
        48,
        [
            ["时间", "CPU / 目标", "CPU 合计", "期望/当前/Ready/实际 Pod"],
            ["11:17:28", "2% / 60%", "1m", "1 / 1 / 1 / 1"],
            ["11:18:05", "472% / 60%", "472m", "3 / 3 / 3 / 3"],
            ["11:18:23", "492% / 60%", "645m", "5 / 5 / 5 / 5"],
            ["11:19:17", "1% / 60%", "505m（终止中）", "1 / 1 / 1 / 5  ← 还不能算缩完"],
            ["11:19:47", "2% / 60%", "2m", "1 / 1 / 1 / 1"],
        ],
        font_size=18,
    )
    d.text((56, 680), "k6：997 请求，9.49 req/s，平均 644.72ms，P95 1124.54ms，错误 0/997。峰值 HPA CPU 501%，商品 Pod CPU 合计 645m。", font=font(20), fill=INK)
    d.text((56, 730), "老师要看全过程：触发 → 响应 → 回落。可以加速剪辑。不要把 HPA 塞进 GitHub Actions。缩容以实际 Pod 回到 1 为准，Terminating 不算完。", font=font(20), fill=RED)
    add(im, "现场放加速录屏。口播 1 到 5 再回到 1。")


def s13():
    im = new_slide()
    d = draw_header(im, "故障隔离：停 product-service，其它业务还在", 13, TOTAL)
    table(
        d,
        (56, 110),
        [420, 1380],
        70,
        [
            ["检查点", "实测（2026-09-01）"],
            ["操作", "实验期间去掉 HPA 后把 product-service 缩到 0；结束 trap 会恢复副本和 HPA"],
            ["商品接口", "HTTP 503，提示「依赖服务暂不可用」"],
            ["订单依赖检查", "HTTP 206，status=degraded，备用提示「商品信息暂不可用，订单查询保持可用」"],
            ["其它服务", "网关 /live 200；用户、订单 /live 均为 alive；三个 Deployment 均为 1/1"],
            ["恢复", "商品接口 200，网关 /ready 200，HPA 重建；无缩成 0 或 HPA 缺失残留"],
        ],
        font_size=20,
    )
    panel(d, (56, 580, 1808, 390))
    d.text((80, 605), "和老师问题的对应", font=font(24), fill=BLUE)
    d.text((80, 660), "压力测试要看到：压力升高后对应微服务 Pod 变多，压力下降后回落。  → 上一页 HPA。", font=font(22), fill=INK)
    d.text((80, 715), "故障处理要看到：一个微服务下线，不严重影响其它微服务的业务；最好返回事先设计好的提示或备用结果。", font=font(22), fill=INK)
    d.text((80, 770), "本组：停的是商品服务，不是把整个系统打死。订单还能查，只是降级。命令 npm run experiment:fault。", font=font(22), fill=INK)
    add(im, "强调 503 和 206 是设计出来的，不是偶然报错。")


def s14():
    im = new_slide()
    d = draw_header(im, "现场演示 4 分钟（可重复真实过程）", 14, TOTAL)
    table(
        d,
        (56, 105),
        [220, 700, 880],
        62,
        [
            ["时间", "做什么", "让老师看见"],
            ["0:00–0:50", "打开 Actions #77，或一次真实 workflow", "提交会跑流水线；红叉不能发版"],
            ["0:50–1:20", "kubectl -n softw-microservices get pods", "网关 + user/product/order 四个负载"],
            ["1:20–2:10", "localhost:8080 走 UC01→UC02→UC04", "登录、搜索、下单支付发货收货"],
            ["2:10–2:30", "Playwright --grep E2E-TC04:", "自动化，不是只靠手点"],
            ["2:30–3:20", "HPA 录屏（加速）", "1 → 3 → 5 → 1"],
            ["3:20–4:00", "停 product-service", "503 / 206，其它 Pod 仍 Running"],
        ],
        font_size=19,
    )
    panel(d, (56, 580, 1808, 390))
    d.text((80, 605), "现场备忘", font=font(24), fill=BLUE)
    lines = [
        "账号 demo-seller@example.com   密码 Demo@123456",
        "前端必须用 http://localhost:8080，不要用 127.0.0.1（CORS）。",
        "公网 Pages 和本机 Docker 不是同一库，不要混着展示「刚才下的单」。",
        "老师要求录屏内嵌 PPT，最好配音和字幕——请把文件放到本页或另插一页视频。",
        "HPA 与故障不要塞进 CI/CD Job，按独立实验脚本跑或放已剪辑录屏。",
    ]
    y = 655
    for t in lines:
        d.text((80, y), "· " + t, font=font(21), fill=INK)
        y += 42
    add(im, "演示按表走。录屏组员嵌入。")


def s15():
    im = new_slide()
    d = draw_header(im, "组内分工 · 大模型 · 过程平台", 15, TOTAL)
    table(
        d,
        (56, 105),
        [280, 1520],
        52,
        [
            ["成员", "主要职责（敏捷记录，权重数字待全组确认）"],
            ["鲁在精", "Scrum Master / 组长：迭代计划、任务协调、风险跟踪、汇报组织"],
            ["浦灵一", "后端重构、安全加固、重构文档"],
            ["王悠然", "前端页面调整、交互优化、构建验证"],
            ["赵紫嫣", "测试计划、接口测试、缺陷跟踪"],
            ["陈子正", "GitHub Actions 流水线、环境配置、部署说明"],
            ["剧博洋", "压测设计、数据库优化、性能分析"],
        ],
        font_size=20,
    )
    panel(d, (56, 500, 1808, 470))
    d.text((80, 525), "大模型怎么用、怎么检查", font=font(24), fill=BLUE)
    y = 580
    for t in [
        "用途：仓库检索、测试建议、文档整理、部署配置检查。",
        "不做的事：不把模型输出直接当成验收结论；不编造未跑过的测试数；不代填权重。",
        "人工检查：读代码、本地命令或 GitHub Actions 真跑、核文档与代码是否同一版本、导出能否打开、凭据扫描。",
        "重要改动走 Pull Request，非作者组员看一眼。",
        "过程只在 GitHub：Issues + Project「软工小学期」+ Actions。没有第二套看板。",
    ]:
        y = draw_wrapped(d, "· " + t, (80, y), font(21), INK, 1740, 32)
        y += 8
    add(im, "权重说待确认。AI 说辅助+人工真跑。")


def s16():
    im = new_slide()
    d = draw_header(im, "演示入口、账号、已知限制", 16, TOTAL)
    table(
        d,
        (56, 110),
        [360, 1440],
        58,
        [
            ["项", "值"],
            ["本地前端", "http://localhost:8080"],
            ["本地健康检查", "http://localhost:3001/api/health"],
            ["公网前端", "https://tchen-0213.github.io/softw/"],
            ["公网后端", "Codespaces softw-defense-demo，3001 端口必须 Public"],
            ["演示账号", "demo-seller@example.com  /  Demo@123456"],
            ["仓库", "https://github.com/tchen-0213/softw"],
            ["看板", "https://github.com/users/tchen-0213/projects/1"],
        ],
        font_size=20,
    )
    panel(d, (56, 620, 1808, 350))
    d.text((80, 645), "限制（主动说，避免被问住）", font=font(24), fill=BLUE)
    y = 700
    for t in [
        "Pages 与本机 Docker 数据不同步。",
        "店铺认证图须为小体积 jpg/png/gif/webp。",
        "CORS 白名单含 localhost:8080，127.0.0.1 可能被拒。",
        "Helm 未使用，Kubernetes 以 YAML 提交。",
    ]:
        d.text((80, y), "· " + t, font=font(22), fill=INK)
        y += 48
    add(im, "限制主动交代一句即可。")


def s17():
    im = new_slide()
    d = draw_header(im, "老师已给的问答（建议原句回答）", 17, TOTAL)
    table(
        d,
        (56, 105),
        [720, 1080],
        88,
        [
            ["问", "答"],
            ["拆成微服务后还要做测试吗？报告交第一阶段还是第二阶段？", "要的，交二阶段的。"],
            ["要演示所有用例吗？", "核心业务场景的用例即可。本组演示 UC01、UC02、UC04。"],
            ["触发扩缩容要演示触发过程吗？要放在 CI/CD 里面吗？", "要从触发到响应的全过程，可适当加速和剪辑，不需要放在 CI/CD 里。"],
            ["压力测试和故障处理具体要演示什么？", "压测看到压力升高后对应微服务 Pod 增加、下降后回落。故障看到一个服务下线不严重影响其它，并返回设计好的提示或备用结果。"],
        ],
        font_size=18,
    )
    panel(d, (56, 580, 1808, 390))
    d.text((80, 605), "本组很可能被追问", font=font(24), fill=BLUE)
    y = 660
    for t in [
        "为什么三个服务不是十二个？→ 身份 / 标的 / 交易过程，不是一个用例一个服务。",
        "评价和聊天为什么在商品服务？→ 围着可交易资源，变化频率接近；不是漏做 evaluation-service。",
        "查询更快为什么 CPU 还更高？→ 快照查询路径短，但网关+双进程+独立库更贵。",
        "测试失败为什么不能部署？→ 流水线 job 依赖，测不过不构建 7 个镜像。",
    ]:
        y = draw_wrapped(d, "· " + t, (80, y), font(21), INK, 1740, 32)
        y += 8
    add(im, "老师原话优先。补充四问各一句。")


def s18():
    im = new_slide()
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, 18, H), fill=BLUE)
    d.rectangle((0, 0, W, 8), fill=BLUE)
    d.text((80, 280), "谢谢老师和助教", font=font(72), fill=BLUE)
    d.text((80, 400), "13组  摸鱼    欢迎提问", font=font(36), fill=INK)
    d.line((80, 470, 420, 470), fill=BLUE2, width=3)
    d.text((80, 510), "https://github.com/tchen-0213/softw", font=font(26), fill=BLUE2)
    d.text((80, 570), "本地  http://localhost:8080", font=font(24), fill=MUTED)
    d.text((80, 620), "重点用例  UC01  注册登录    UC02  搜索    UC04  下单履约", font=font(24), fill=INK)
    paste_contain(im, PNG / "33-MS-SPLIT.png", (980, 180, 860, 720))
    add(im, "停住等提问。")


BUILDERS = [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11, s12, s13, s14, s15, s16, s17, s18]


def pack_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i, im in enumerate(SLIDES):
        path = OUT_DIR / f"S{i + 1:02d}.png"
        im.save(path, "PNG", optimize=True)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height)
        notes = slide.notes_slide
        notes.notes_text_frame.text = NOTES[i]
    prs.save(PPTX)


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for fn in BUILDERS:
        fn()
    if len(SLIDES) != TOTAL:
        print(f"expected {TOTAL} slides, got {len(SLIDES)}", file=sys.stderr)
        return 1
    pack_pptx()
    print(PPTX)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
